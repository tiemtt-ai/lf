"""Small application test inputs; requires python-docx/python-pptx/openpyxl.
Convert to legacy formats with LibreOffice in a disposable user profile.
"""
from pathlib import Path
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
out = Path(__file__).resolve().parent
word = Document()
word.add_heading('LearnForge Office fixture', level=1)
word.add_paragraph('DOCUMENT_ACCEPTANCE_ALPHA')
word.save(out / 'office.docx')
slides = Presentation()
slide = slides.slides.add_slide(slides.slide_layouts[1])
slide.shapes.title.text = 'LearnForge Office fixture'
slide.placeholders[1].text = 'DOCUMENT_ACCEPTANCE_ALPHA'
slides.save(out / 'office.pptx')
book = Workbook()
book.active.title = 'Sales'
book.active.append(['DOCUMENT_ACCEPTANCE_ALPHA', 'Count'])
book.active.append(['ALPHA', 10])
book.create_sheet('Notes').append(['Second worksheet'])
book.save(out / 'office.xlsx')
